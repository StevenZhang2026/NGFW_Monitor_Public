# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

TEMPLATE = '/Users/szhang/Documents/claude/NGFW_bytedance/docs/2026 PANW Corporate Presentation Template.pptx'
OUTPUT = '/Users/szhang/Documents/claude/NGFW_bytedance/docs/NGFW_Monitor_Vibecoding.pptx'

prs = Presentation(TEMPLATE)

# Delete all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Layout indices
LY_TITLE = 0          # TITLE: CENTER_TITLE(0) + SUBTITLE(1) + SUBTITLE(2)
LY_SECTION = 10       # SECTION_HEADER: TITLE(0)
LY_BODY = 13          # TITLE_AND_BODY_1_1: TITLE(0) + BODY(1)
LY_TWO_COL = 11       # TITLE_AND_TWO_COLUMNS_3: BODY(1) + TITLE(0) + BODY(2)


def add_title_slide(title, subtitle, subtitle2=''):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE])
    slide.placeholders[0].text = title
    slide.placeholders[1].text = subtitle
    if subtitle2:
        slide.placeholders[2].text = subtitle2
    return slide


def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION])
    slide.placeholders[0].text = title
    return slide


def add_body_slide(title, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_BODY])
    slide.placeholders[0].text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    lines = body_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        # Indent bullet points
        if line.startswith('  '):
            p.level = 1
        elif line.startswith('    '):
            p.level = 2
    return slide


def add_two_col_slide(title, left_text, right_text):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TWO_COL])
    slide.placeholders[0].text = title
    # Left column
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(left_text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        if line.startswith('  '):
            p.level = 1
    # Right column
    tf = slide.placeholders[2].text_frame
    tf.clear()
    for i, line in enumerate(right_text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        if line.startswith('  '):
            p.level = 1
    return slide


# ============================================================
# SLIDES
# ============================================================

# 1. Title
add_title_slide(
    'NGFW Monitor',
    'Palo Alto 防火墙集中监控与分析平台',
    '从 0 到 24000 行代码的 Vibecoding 之旅'
)

# 2. Section: Why
add_section_slide('为什么要做这个项目？')

# 3. Pain points
add_body_slide('痛点', '''
20 多台 PA 防火墙，每台都要单独登录看状态
CPU 飙了没人知道，直到用户投诉"网怎么这么慢"
领导问"上周威胁怎么样"，得手动截图拼 Excel
新人来了："这 20 台设备我该看哪几台？"
告警靠人肉巡检，等发现问题时已经晚了
''')

# 4. Vision
add_body_slide('我们想要什么', '''
一个平台，所有设备状态一目了然
告警自动推送到飞书/邮件，不用人肉巡检
报表自动发领导邮箱，不用手动截图
新人一句话就能查到想要的数据
笔记本能跑，服务器也能跑，代码一套搞定
''')

# 5. Section: Tech
add_section_slide('技术选型')

# 6. Tech stack
add_two_col_slide('技术栈一览',
'''后端
Python 3.11 + FastAPI
Celery + Redis（定时采集）
SQLAlchemy + asyncpg
PAN-OS XML API + SSH CLI''',
'''前端 & 数据
React 18 + TypeScript
ECharts + Ant Design
PostgreSQL + TimescaleDB
Docker Compose 部署''')

# 7. Section: Architecture
add_section_slide('系统架构')

# 8. Architecture
add_body_slide('整体架构', '''
React 前端 → Nginx (SSL) → FastAPI → TimescaleDB
Celery Beat → Worker → PA 防火墙 (API/SSH)
Worker → 告警引擎 → 飞书 / 邮件
FastAPI → LLM 服务（仅意图解析，数据不出门）
全部容器化，Docker Compose 一条命令启动
笔记本 ~1.2GB RAM | 服务器横向扩展
''')

# 9. Section: Day 1
add_section_slide('Day 1 — "先聊别写码"')

# 10. Day 1 Discussion
add_body_slide('架构讨论（1 小时）', '''
用户第一句话："先不要写代码，先聊情况架构。"

讨论结果：
采集频率谁说了算？→ "人改可以，系统不自作主张"
数据会降采样吗？→ "原始数据全量保留，展示粒度用户选"
资源不够怎么办？→ "告警通知，不准静默降级"
指标能扩展吗？→ "80% 配置驱动，20% 写插件"
部署到哪里？→ "笔记本 Docker 先跑，后期迁服务器"

讨论完 → 52 个文件一口气生成
然后用户说："好的再见，等我找到防火墙再继续"
''')

# 11. Day 1 PA-440
add_body_slide('PA-440 接入实测', '''
用户回来了，带着一台 PA-440（192.168.1.254）

踩坑四连：
1. Docker 镜像编译了 46 分钟 → 精简依赖，3 分钟搞定
2. bcrypt 版本冲突 → 直接用 bcrypt，甩掉 passlib
3. 容器连不上防火墙！→ Docker Desktop "Access local network"
4. Collector 没注册 → worker 启动时忘了 import

最终结果：CPU 6.1%、内存 28.57%、会话 0/199998
端到端采集成功！
''')

# 12. Section: Core Features
add_section_slide('核心能力')

# 13. Monitoring
add_body_slide('实时监控', '''
12 个内置指标，开箱即用：
  CPU / 内存 / 会话数 / Packet Descriptor
  接口吞吐 / 温度 / HA 状态
  应用 Top10 / 威胁 Top10

Dashboard 四宫格实时展示
PAN-OS XML API + SSH CLI 双通道采集
每设备单连接复用，不浪费管理面资源
设备状态自动检测：采集成功→online，失败→offline
''')

# 14. Custom metrics
add_body_slide('为什么指标要可扩展？', '''
用户说："监控项目或指标需要可以后期增加。"

设计方案：
80% 场景：Web UI 填 CLI 命令 + 正则 → 新指标上线
20% 场景：写 Collector 插件（继承基类 + 装饰器注册）
指标定义存数据库，热更新，不用重启服务

支持解析器：xpath / regex / regex_cdata / xpath_multi

"想监控什么就监控什么，不用求开发改代码"
''')

# 15. Groups & Scope
add_body_slide('为什么加分组和权限隔离？', '''
场景：团队 20 台防火墙分布在 3 个数据中心
小王只管北京的 5 台，不该看到上海的数据

RBAC + Scope 双层控制：

RBAC 管"能做什么"：
  admin → 全能 | operator → 操作 | viewer → 只读

Scope 管"能看什么"：
  按设备分组授权，每人只看自己组的数据
  无 scope 分配 = 全局访问（向后兼容）

"你的防火墙你负责，别人的数据你看不到"
''')

# 16. Alerts
add_body_slide('告警体系', '''
三种告警类型：
  阈值告警 — CPU > 80%？立刻通知
  异常检测 — 会话数飙升 3 个标准差？可疑
  趋势预测 — 按当前速度 12 周后 CPU 爆了

通知渠道：飞书群机器人 + SMTP 邮件

踩坑记录（飞书通知调通花了 3 轮）：
  1. HTTP 200 不等于成功，要检查 body.code
  2. 企业网有 SSL inspection → verify=False
  3. 机器人关键词"防火墙"→ 消息模板必须带

v2.1 新增：通知冷却（别 1 分钟发 60 条把人逼疯）
''')

# 17. Reports
add_body_slide('为什么加报表模块？', '''
领导的灵魂拷问：
  "上周防火墙安全态势怎么样？"
  "CPU 趋势是不是在涨？还能撑多久？"

领导不会登录你的 Dashboard。
领导只看邮箱里的 PDF。

解决方案 — 自动报表系统：
  周报（每周一 08:00）/ 月报（每月 1 日）
  numpy 线性回归计算趋势 + 预测容量
  matplotlib 图表 → weasyprint → PDF
  aiosmtplib 发邮件，PDF 附件直达领导邮箱

模板化结论（不用 LLM）：
  "CPU 呈上升趋势，预计 12 周后触及告警阈值"
''')

# 18. AI Copilot
add_body_slide('为什么加 AI Copilot？', '''
用户说："能自然语言查询就好了"
  比如"给我 3 天内威胁 Top 10"

关键顾虑："数据不会传到云端大模型吗？"

混合架构 — 数据不出门：
  1. 用户提问 → 只把问题文本发给 LLM
  2. LLM 返回结构化意图 {action, params}
  3. 后端内部执行 SQL 查询（数据不出门！）
  4. 模板格式化为 Markdown 表格返回前端

模型可配置：管理员设置 API Base + Key + Model
支持 DeepSeek / OpenAI / Ollama（本地更安全）
''')

# 19. Section: Bugs
add_section_slide('那些有趣的 Bug')

# 20. Bug stories
add_body_slide('调试趣事', '''
ACC 趋势图 tooltip 数值错位：
  根因：不同应用有数据的时间戳不一样
  修复：所有系列对齐到相同时间戳，缺失填 0

Copilot 永远说"无法理解您的问题"：
  根因：parse_intent 遇到错误静默 return None
  实际错误：企业代理返回了 503 HTML 登录页
  修复：加 IntentError + 具体错误信息暴露给用户

Docker 容器连不上局域网防火墙：
  根因：macOS Docker Desktop VM 与宿主网络隔离
  修复：Settings → Access local network
  教训：这个选项默认关闭，文档里没写...
''')

# 21. Section: Principles
add_section_slide('设计原则')

# 22. Principles
add_body_slide('五条铁律', '''
1. 采集频率由管理员控制
   系统永远不自动降频，人改可以，系统不自作主张

2. 原始数据全量保留
   不丢点不合并不降采样，展示粒度用户选

3. 指标可扩展
   80% 配置零代码，20% 写插件

4. 环境无关可迁移
   同一套代码 .env + overlay 适配笔记本/服务器

5. 资源不足时明确告警
   不静默降级，通知用户决定扩容还是调整
''')

# 23. Section: Stats
add_section_slide('数字说话')

# 24. Stats
add_two_col_slide('项目统计',
'''开发周期
  4 天（8 个会话）

总代码量
  ~24,000 行
  Python + TypeScript + YAML

版本迭代
  v1.0 → v1.1 → v1.2 → v2.0 → v2.1''',
'''功能模块
  采集 · 展示 · 告警
  权限 · 报表 · AI 助手
  部署工具

已验证设备
  PA-440 (PAN-OS 10.1)
  兼容 PA-5500 / PA-7000''')

# 25. Timeline
add_body_slide('版本演进时间线', '''
Day 1 · v1.0 — 从需求讨论到端到端采集成功
  52 文件骨架 + PA-440 实测验证

Day 2 · v1.1 — 权限体系 + ACC 数据 + 安全加固
  设备分组 / Scope 隔离 / SQL注入修复 / 飞书通知

Day 3 · v1.2 — ACC 重构 + 安装工具
  Log Query API / install.sh 一键部署

Day 3 · v2.0 — 报表模块
  PDF 周报月报 / numpy 趋势预测 / 邮件推送

Day 4 · v2.1 — AI Copilot + 告警优化
  自然语言查询 / 通知冷却 / 图表修复 / 架构图
''')

# 26. Section: Insights
add_section_slide('Vibecoding 心得')

# 27. Insights
add_body_slide('协作模式总结', '''
先讨论再动手
  1 小时的架构讨论省了 1 周的返工

用真设备验证
  代码生成容易，能跑通才是真功夫

每个功能都有"为什么"
  不是技术炫技，是解决真实痛点

迭代而非一步到位
  先跑通 → 再完善 → 再优化

AI 是协作者不是替代品
  人定方向和约束，AI 负责实现和排错
''')

# 28. End slide
add_title_slide(
    'Thank You',
    'github.com/StevenZhang2026/NGFW_Monitor_Public',
    'Powered by Vibecoding — Human + AI Collaboration'
)

# Save
prs.save(OUTPUT)
print(f'Done: {OUTPUT}')
print(f'Slides: {len(prs.slides)}')
